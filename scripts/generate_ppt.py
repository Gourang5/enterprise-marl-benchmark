"""
Generate Enterprise MARL Benchmark presentation as a .pptx file.
Run: python scripts/generate_ppt.py
Output: deliverables/enterprise_marl_benchmark.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT     = RGBColor(0x2E, 0x9C, 0xFF)
ACCENT2    = RGBColor(0x00, 0xD4, 0x9A)
WARN       = RGBColor(0xFF, 0xA5, 0x2A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF4, 0xF8)
MID_GREY   = RGBColor(0x8A, 0x9B, 0xAD)
DARK_TEXT  = RGBColor(0x1A, 0x2B, 0x3C)
CARD_BG    = RGBColor(0xF7, 0xFA, 0xFC)
GREEN_OK   = RGBColor(0x1A, 0xB0, 0x70)
RED_NO     = RGBColor(0xE5, 0x3E, 0x3E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# Helpers

def add_rect(slide, l, t, w, h, fill=None, line=None, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.width = Pt(line)
        shape.line.color.rgb = line_color or MID_GREY
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, valign=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def dark_slide(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK_BG)


def section_bar(slide, label):
    add_rect(slide, 0, 0, 13.33, 0.55, fill=ACCENT)
    add_text(slide, label, 0.3, 0.08, 12, 0.4, size=13, bold=True, color=WHITE)


def slide_title(slide, title, subtitle=None, dark=False):
    tc = WHITE if dark else DARK_TEXT
    add_text(slide, title, 0.5, 0.6, 12.3, 0.8, size=28, bold=True, color=tc)
    if subtitle:
        add_text(slide, subtitle, 0.5, 1.35, 12.3, 0.5,
                 size=15, color=ACCENT, italic=True)


def chip(slide, label, l, t, w=1.6, h=0.32, bg=ACCENT, tc=WHITE):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, label, l+0.08, t+0.04, w-0.16, h-0.08,
             size=10, bold=True, color=tc, align=PP_ALIGN.CENTER)


# SLIDE 1 - Title / Hero
sl = prs.slides.add_slide(BLANK)
dark_slide(sl)
add_rect(sl, 0, 6.8, 13.33, 0.7, fill=ACCENT)

add_text(sl, "Enterprise Multi-Agent RL Benchmark",
         0.7, 1.2, 11.9, 1.1, size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl,
    "Can autonomous agents coordinate across applications, roles, and permissions\n"
    "to complete realistic enterprise workflows?",
    1.0, 2.4, 11.3, 1.0,
    size=17, color=RGBColor(0xB0, 0xD8, 0xFF), align=PP_ALIGN.CENTER, italic=True)

stats = [
    ("5", "Heterogeneous\nAgents"),
    ("5", "Simulated\nApps"),
    ("6", "Dependency-\nGated Tasks"),
    ("118", "Automated\nTests"),
]
for i, (num, lbl) in enumerate(stats):
    x = 1.3 + i * 2.75
    add_rect(sl, x, 3.7, 2.3, 1.5, fill=RGBColor(0x1A, 0x3A, 0x5C))
    add_text(sl, num, x, 3.75, 2.3, 0.8, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, x, 4.45, 2.3, 0.65, size=11, color=MID_GREY, align=PP_ALIGN.CENTER)

add_text(sl, "Shared SQLite State  ·  Deterministic Evaluator  ·  Fully Reproducible from Seed",
         0.5, 5.6, 12.3, 0.4, size=12, color=MID_GREY, align=PP_ALIGN.CENTER)
add_text(sl, "v1.3  |  2025", 0.3, 6.88, 4, 0.3, size=11, color=WHITE)


# SLIDE 2 - Gap / Comparison Table
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Why Existing Benchmarks Miss Enterprise Work",
            subtitle="Enterprise coordination demands all these dimensions simultaneously")

headers = ["Dimension", "Typical Benchmark", "This Benchmark"]
rows = [
    ("Agents",        "1",                          "5, heterogeneous"),
    ("Apps / state",  "1",                          "5, shared persistent DB"),
    ("Permissions",   "None",                       "Per-role x per-app x per-sheet"),
    ("Horizon",       "1-5 steps",                  "8-17 subgoals, 25-45 steps"),
    ("Dependencies",  "None",                       "Gated DAG: later credit blocked"),
    ("Information",   "Symmetric",                  "Asymmetric per-agent observations"),
    ("Evaluation",    "String match / LLM judge",   "Deterministic DB state check"),
]
col_w = [3.0, 3.8, 4.5]
col_x = [0.5, 3.7, 7.7]
row_h = 0.56
start_y = 2.05

for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    bg = DARK_BG if ci > 0 else RGBColor(0x3A, 0x5A, 0x7A)
    add_rect(sl, cx, start_y, cw, row_h, fill=bg)
    add_text(sl, hdr, cx+0.1, start_y+0.12, cw-0.2, row_h-0.15,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    y = start_y + (ri+1)*row_h
    for ci, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        if ci == 0:
            bg, tc = CARD_BG, DARK_TEXT
        elif ci == 1:
            bg, tc = RGBColor(0xE8, 0xEC, 0xF0), RGBColor(0x6A, 0x6A, 0x6A)
        else:
            bg, tc = RGBColor(0xE6, 0xF6, 0xEF), RGBColor(0x0A, 0x6A, 0x42)
        add_rect(sl, cx, y, cw, row_h, fill=bg, line=0.5, line_color=RGBColor(0xD0, 0xD8, 0xE0))
        add_text(sl, val, cx+0.12, y+0.1, cw-0.24, row_h-0.15, size=11, color=tc, bold=(ci==2))

add_rect(sl, 0.5, 6.45, 12.33, 0.5, fill=RGBColor(0x0D, 0x4F, 0x8C))
add_text(sl,
    "Design thesis: a single evaluation environment that combines all of these dimensions simultaneously.",
    0.7, 6.5, 12.0, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# SLIDE 3 - Environment Overview
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Environment Overview",
            subtitle="One synthetic company · five app simulators · one shared database")

apps = ["Gmail", "Slack", "Jira", "Calendar", "Sheets"]
app_colors = [
    RGBColor(0xEA, 0x43, 0x35), RGBColor(0x4A, 0x15, 0x4B),
    RGBColor(0x00, 0x52, 0xCC), RGBColor(0x0F, 0x9D, 0x58), RGBColor(0x34, 0xA8, 0x53),
]
for i, (app, color) in enumerate(zip(apps, app_colors)):
    x = 1.3 + i * 2.15
    add_rect(sl, x, 1.85, 1.8, 0.7, fill=color)
    add_text(sl, app, x, 1.9, 1.8, 0.6, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "↓", 6.4, 2.62, 0.5, 0.4, size=22, color=ACCENT, align=PP_ALIGN.CENTER)
add_rect(sl, 3.5, 3.05, 6.3, 0.7, fill=DARK_BG)
add_text(sl, "SQLite: Shared Company State", 3.5, 3.1, 6.3, 0.6,
         size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(sl, 0.8, 4.0, 4.5, 1.1, fill=RGBColor(0xE6, 0xF0, 0xFF))
add_text(sl, "Agent Observations\n(role-filtered per agent)", 0.8, 4.05, 4.5, 1.0,
         size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_rect(sl, 8.0, 4.0, 4.5, 1.1, fill=RGBColor(0xFF, 0xF3, 0xE0))
add_text(sl, "Deterministic Evaluator\n(private, state-based)", 8.0, 4.05, 4.5, 1.0,
         size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)

bullets = [
    "Every episode resets from a deterministic integer seed",
    'Evaluator state is NOT visible to agents - info["eval"] is separate',
    "Semantic tool actions: malformed calls are invalid transitions, not crashes",
    "An agent cannot succeed by claiming 'task complete' in Slack - the DB must match",
]
for i, b in enumerate(bullets):
    add_text(sl, f"•  {b}", 0.6, 5.3 + i*0.38, 12.1, 0.35, size=12, color=DARK_TEXT)


# SLIDE 4 - Architecture
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Architecture",
            subtitle="Gymnasium-compatible RL interface with semantic action dispatch")

pipeline = [
    (ACCENT,                     "Policy / Controller",               "Issues Action(agent, app, action_type, params)"),
    (RGBColor(0x1A, 0x5C, 0x8C), "EnterpriseEnv",                    "Role-filtered obs · Action validator · Reward engine"),
    (RGBColor(0x0A, 0x7A, 0x5C), "App Simulators (thin, stateless)", "Pure DB reads/writes: Gmail · Slack · Jira · Calendar · Sheets"),
    (DARK_BG,                    "SQLite Shared State",               "Single source of truth across all apps"),
    (RGBColor(0x6A, 0x2A, 0x8C), "Private Verifier",                 'DAG + subgoal checks → reward + info["eval"]'),
]
for i, (color, label, desc) in enumerate(pipeline):
    y = 1.85 + i * 0.95
    add_rect(sl, 1.5, y, 10.3, 0.72, fill=color)
    add_text(sl, label, 1.7, y+0.05, 3.8, 0.4, size=13, bold=True, color=WHITE)
    add_text(sl, desc,  5.6, y+0.08, 6.0, 0.4, size=12, color=RGBColor(0xC8, 0xE0, 0xFF))
    if i < len(pipeline)-1:
        add_text(sl, "↓", 6.6, y+0.72, 0.4, 0.25, size=14, color=ACCENT, align=PP_ALIGN.CENTER)

notes = [
    "Execution: turn-based, actor-selected sequential (not truly parallel)",
    "RL interface: Gymnasium-compatible  reset(seed) / step(action) / close()",
    "PettingZoo AEC adapter available - EXPERIMENTAL (not fully tested)",
]
for i, n in enumerate(notes):
    add_text(sl, f"•  {n}", 0.7, 6.55 + i*0.29, 12.0, 0.28, size=11,
             color=MID_GREY, italic=(i==2))


# SLIDE 5 - Multi-Agent Model
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Multi-Agent Model",
            subtitle="Five agents · different roles · different app access · different information")

headers5 = ["Agent", "Role", "Gmail", "Slack", "Jira", "Calendar", "Sheets"]
# No-access shown as "-" (regular hyphen, not em dash)
agents_data = [
    ("pm_01",      "Project Manager",  "Yes", "Yes", "Yes", "Yes", "Owner R/W", RGBColor(0x1A, 0x5C, 0xFF)),
    ("eng_01",     "Engineer",         "-",   "Yes", "Yes", "Yes", "Viewer",    RGBColor(0x1A, 0xB0, 0x70)),
    ("product_01", "Product Manager",  "-",   "Yes", "Yes", "-",   "Editor",    RGBColor(0xFF, 0x8C, 0x1A)),
    ("mgr_01",     "Eng. Manager",     "-",   "Yes", "Yes", "Yes", "-",         RGBColor(0xE5, 0x3E, 0x3E)),
    ("cs_01",      "Customer Success", "Yes", "Yes", "Yes", "-",   "-",         RGBColor(0x8A, 0x4A, 0xCC)),
]

col_widths = [1.6, 2.5, 0.95, 0.95, 0.85, 1.1, 1.43]
col_starts = [0.4]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

row_h = 0.72
hy = 2.0

# Header row
for ci, (hdr, cx, cw) in enumerate(zip(headers5, col_starts, col_widths)):
    add_rect(sl, cx, hy, cw-0.05, row_h, fill=DARK_BG)
    add_text(sl, hdr, cx+0.05, hy, cw-0.1, row_h,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)

# Data rows
for ri, (ag_id, role, *cells, color) in enumerate(agents_data):
    y = hy + (ri+1)*row_h
    row_vals = [ag_id, role] + cells
    for ci, (val, cx, cw) in enumerate(zip(row_vals, col_starts, col_widths)):
        bg = CARD_BG if ri%2==0 else RGBColor(0xEF, 0xF3, 0xF7)
        add_rect(sl, cx, y, cw-0.05, row_h, fill=bg,
                 line=0.4, line_color=RGBColor(0xD0, 0xD8, 0xE2))
        if ci == 0:
            tc = color
        elif val == "Yes":
            tc = GREEN_OK
        elif val == "-":
            tc = MID_GREY
        else:
            tc = DARK_TEXT
        add_text(sl, val, cx+0.05, y, cw-0.1, row_h,
                 size=11, bold=(ci==0), color=tc, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)

add_text(sl, "Who can act: per-agent tool whitelist enforced on every call",
         0.5, 6.65, 12.3, 0.35, size=12, color=DARK_TEXT)
add_text(sl, "What each agent sees: only their own inbox, accessible channels, permitted sheets",
         0.5, 7.0, 12.3, 0.35, size=12, color=DARK_TEXT)


# SLIDE 6 - Vendor Onboarding Task
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Long-Horizon Task: Vendor Onboarding",
            subtitle="8 subgoals · 40-step horizon · 5 apps · 4 roles must act")

steps = [
    ("pm_01",      "Discover vendor email + find procurement ticket",  "Gmail → Jira",  RGBColor(0x1A, 0x5C, 0xFF)),
    ("product_01", "Clear legal contract ticket",                       "Jira",          RGBColor(0xFF, 0x8C, 0x1A)),
    ("eng_01",     "Confirm IT provisioning ticket",                    "Jira",          RGBColor(0x1A, 0xB0, 0x70)),
    ("mgr_01",     "Approve procurement ticket",                        "Jira",          RGBColor(0xE5, 0x3E, 0x3E)),
    ("pm_01",      "Mark vendor ACTIVE in tracker sheet",               "Sheets",        RGBColor(0x1A, 0x5C, 0xFF)),
    ("pm_01",      "Schedule kickoff meeting",                          "Calendar",      RGBColor(0x1A, 0x5C, 0xFF)),
    ("pm_01",      "Announce completion in Slack",                      "Slack",         RGBColor(0x1A, 0x5C, 0xFF)),
]

for i, (agent, action, app, color) in enumerate(steps):
    y = 1.85 + i * 0.68
    add_rect(sl, 0.5, y, 1.5, 0.55, fill=color)
    add_text(sl, agent, 0.5, y+0.1, 1.5, 0.38, size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 2.15, y, 8.3, 0.55, fill=CARD_BG,
             line=0.5, line_color=RGBColor(0xC8, 0xD8, 0xE8))
    add_text(sl, action, 2.3, y+0.1, 8.0, 0.38, size=12, color=DARK_TEXT)
    add_rect(sl, 10.6, y, 2.2, 0.55, fill=RGBColor(0xE8, 0xF4, 0xFF))
    add_text(sl, app, 10.65, y+0.1, 2.1, 0.38, size=11,
             color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    if i in (1, 2):
        add_text(sl, "parallel branch", 0.5, y+0.57, 2.5, 0.22,
                 size=8, color=MID_GREY, italic=True)

props = ["Parallel branches (legal + IT)", "Strict RBAC on sheet write",
         "Negation-hardened verification", "Role separation + info asymmetry"]
for i, p in enumerate(props):
    chip(sl, p, 0.5 + i*3.2, 6.9, w=3.0, h=0.36, bg=DARK_BG, tc=WHITE)


# SLIDE 7 - Task Dependency DAG
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Task Dependency DAG: Vendor Onboarding",
            subtitle="Credit fires when DB state matches the predicate, not when an agent says so")

nodes = [
    ("discover_request",    "pm_01",      3.5, 1.65, RGBColor(0x1A, 0x5C, 0xFF)),
    ("find_main_ticket",    "pm_01",      3.5, 2.55, RGBColor(0x1A, 0x5C, 0xFF)),
    ("legal_review",        "product_01", 1.0, 3.5,  RGBColor(0xFF, 0x8C, 0x1A)),
    ("it_provisioning",     "eng_01",     6.0, 3.5,  RGBColor(0x1A, 0xB0, 0x70)),
    ("manager_approval",    "mgr_01",     3.5, 4.45, RGBColor(0xE5, 0x3E, 0x3E)),
    ("update_vendor_sheet", "pm_01",      3.5, 5.4,  RGBColor(0x1A, 0x5C, 0xFF)),
    ("schedule_kickoff",    "pm_01",      1.0, 6.35, RGBColor(0x1A, 0x5C, 0xFF)),
    ("announce_live",       "pm_01",      6.0, 6.35, RGBColor(0x1A, 0x5C, 0xFF)),
]
for (label, agent, nx, ny, color) in nodes:
    add_rect(sl, nx, ny, 3.0, 0.65, fill=color)
    add_text(sl, label, nx+0.1, ny+0.04, 2.8, 0.3,  size=11, bold=True, color=WHITE)
    add_text(sl, agent, nx+0.1, ny+0.33, 2.8, 0.22, size=9, color=RGBColor(0xC0, 0xD8, 0xFF))

for (ax, ay, sym) in [(5.0,2.3,"↓"),(5.0,3.25,"↓"),(5.0,4.2,"↓"),(5.0,5.15,"↓"),
                       (2.5,4.2,"↙"),(7.5,4.2,"↘"),(2.5,6.1,"↓"),(7.5,6.1,"↓")]:
    add_text(sl, sym, ax, ay, 0.4, 0.3, size=16, color=ACCENT, align=PP_ALIGN.CENTER)

for i, b in enumerate([
    "Subgoals are locked until all prerequisites complete",
    "Parallel branches (legal + IT) both required before schedule_kickoff",
    'Verifier is negation-hardened: "NOT provisioned" fails · "Provisioned and confirmed" passes',
]):
    add_text(sl, f"•  {b}", 9.5, 2.8 + i*0.55, 3.6, 0.48, size=10, color=DARK_TEXT)


# SLIDE 8 - Partial Observability
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Partial Observability & Permissions",
            subtitle="Agents must search and discover - ticket IDs are never pre-known")

add_rect(sl, 0.5, 2.0, 5.5, 3.8, fill=RGBColor(0xE6, 0xF0, 0xFF), line=1.0, line_color=ACCENT)
add_text(sl, "What each agent CAN see", 0.6, 2.05, 5.3, 0.45, size=13, bold=True, color=ACCENT)
for i, item in enumerate([
    "Own email inbox headers (body requires read_email)",
    "Slack channels they belong to",
    "Jira tickets returned by their own search query",
    "Own calendar events only",
]):
    add_text(sl, f"Yes  {item}", 0.7, 2.58 + i*0.58, 5.1, 0.5, size=12, color=GREEN_OK)

add_rect(sl, 7.3, 2.0, 5.5, 3.8, fill=RGBColor(0xFF, 0xF0, 0xE8), line=1.0, line_color=RED_NO)
add_text(sl, "What NO agent can see", 7.4, 2.05, 5.3, 0.45, size=13, bold=True, color=RED_NO)
for i, item in enumerate([
    "Other agents' inboxes",
    "Verifier progress or subgoal status",
    "Full Jira database (search required)",
    "Sheet cells outside membership role",
]):
    add_text(sl, f"No   {item}", 7.45, 2.58 + i*0.58, 5.1, 0.5, size=12, color=RED_NO)

add_rect(sl, 6.16, 2.4, 1.0, 3.0, fill=LIGHT_GREY)
add_text(sl, "🔒", 6.2, 3.6, 0.9, 0.8, size=26, align=PP_ALIGN.CENTER)

add_rect(sl, 0.5, 6.1, 12.3, 0.8, fill=DARK_BG)
add_text(sl,
    'Wrong permission → {"success": false, "message": "Permission denied"} - not silent failure',
    0.7, 6.2, 12.0, 0.5, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# SLIDE 9 - Reward Design
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Reward Design & Deterministic Verification",
            subtitle="Shaped signal across subgoal DAG · terminal success · anti-hacking penalties")

rewards = [
    ("Valid tool action",   "+0.25",   "Execution credit",              GREEN_OK),
    ("Subgoal unlocked",    "+8x(1/N)","Shaped, per new subgoal",       ACCENT),
    ("Coordination bonus",  "+2.0",    "Multi-agent handoff required",  ACCENT2),
    ("Terminal success",    "+75.0",   "Episode completion",            GREEN_OK),
    ("Redundant action",    "-1.0",    "Penalises reward farming",      RED_NO),
    ("Step cost",           "-0.10",   "Efficiency pressure",           RGBColor(0xCC, 0x60, 0x00)),
    ("Timeout",             "-15.0",   "Horizon enforcement",           RED_NO),
]
col_x2 = [0.5, 3.5, 5.5]
col_w2 = [3.0, 2.0, 7.3]
hy2 = 2.0
add_rect(sl, 0.5, hy2, 12.3, 0.5, fill=DARK_BG)
for hdr, cx, cw in zip(["Signal", "Value", "Purpose"], col_x2, col_w2):
    add_text(sl, hdr, cx+0.05, hy2+0.08, cw, 0.35, size=11, bold=True, color=WHITE)

for ri, (sig, val, purp, color) in enumerate(rewards):
    y = hy2 + 0.5 + ri*0.52
    bg = CARD_BG if ri%2==0 else RGBColor(0xEF, 0xF3, 0xF7)
    add_rect(sl, 0.5, y, 12.3, 0.49, fill=bg, line=0.4, line_color=RGBColor(0xD0, 0xD8, 0xE0))
    add_text(sl, sig,  0.6,  y+0.09, 2.85, 0.32, size=11, color=DARK_TEXT)
    add_text(sl, val,  3.55, y+0.09, 1.9,  0.32, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, purp, 5.55, y+0.09, 7.2,  0.32, size=11, color=DARK_TEXT)

add_rect(sl, 0.5, 6.6, 12.3, 0.65, fill=RGBColor(0xE6, 0xF6, 0xEF))
add_text(sl,
    'Verification: Success = state predicate over SQLite - not string match, not LLM judge. '
    'Clause-boundary negation detection: "NOT approved" → fail · "Approved. No further action." → pass',
    0.7, 6.68, 12.0, 0.5, size=11, color=RGBColor(0x0A, 0x5A, 0x36))


# SLIDE 10 - Evaluation & Baselines (Random removed)
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Evaluation, Baselines & Replay",
            subtitle="Three-tier policy taxonomy · Wilson 95% CI · HTML trajectory viewer")

# Bar chart - 3 bars (Random removed), wider and centred
bars = [
    ("Zero-Shot LLM",   0,   WARN,     "0 / 30"),
    ("Hint-Guided LLM", 80,  ACCENT,   "24 / 30"),
    ("Oracle",          100, GREEN_OK, "30 / 30"),
]
bar_area_x = 1.0
bar_area_y = 1.85
bar_w      = 3.2
bar_gap    = 0.55
bar_max_h  = 2.8

for i, (label, pct, color, result) in enumerate(bars):
    x = bar_area_x + i*(bar_w + bar_gap)
    h = bar_max_h * pct / 100 if pct > 0 else 0.12
    y = bar_area_y + bar_max_h - h
    add_rect(sl, x, y, bar_w, h, fill=color)
    add_text(sl, result, x, y-0.4, bar_w, 0.35,
             size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, label, x, bar_area_y+bar_max_h+0.08, bar_w, 0.45,
             size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_rect(sl, 0.9, bar_area_y+bar_max_h, 11.5, 0.04, fill=DARK_BG)

# Table - 3 rows
table_y = 5.1
headers3 = ["Policy", "Result", "Notes"]
col_x3 = [0.5, 3.8, 6.0]
col_w3 = [3.3, 2.2, 7.1]
add_rect(sl, 0.5, table_y, 12.3, 0.45, fill=DARK_BG)
for hdr, cx, cw in zip(headers3, col_x3, col_w3):
    add_text(sl, hdr, cx+0.05, table_y+0.08, cw, 0.3, size=11, bold=True, color=WHITE)

table_rows = [
    ("Zero-Shot LLM",   "0/30 ep (5 ep/task)",   "qwen2.5:3b - model-dependent result"),
    ("Hint-Guided LLM", "24/30 ep (5 ep/task)",  "SOP debug baseline - NOT autonomous capability"),
    ("Oracle",          "30/30 episodes",         "Scripted upper bound - proves solvability"),
]
for ri, row in enumerate(table_rows):
    y = table_y + 0.45 + ri*0.46
    bg = CARD_BG if ri%2==0 else RGBColor(0xEF, 0xF3, 0xF7)
    add_rect(sl, 0.5, y, 12.3, 0.43, fill=bg, line=0.3, line_color=RGBColor(0xD0, 0xD8, 0xE0))
    for val, cx, cw in zip(row, col_x3, col_w3):
        add_text(sl, val, cx+0.05, y+0.08, cw, 0.28, size=10, color=DARK_TEXT)


# SLIDE 11 - Research Use Cases & Limitations
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Research Use Cases & Honest Limitations",
            subtitle="What the benchmark measures · where it stops")

capabilities = [
    "Planning: multi-step reasoning across apps and roles",
    "Tool use: structured semantic action under partial observability",
    "Coordination: one agent's output blocks another's subgoal",
    "Delegation: transferring work across roles at the correct moment",
    "Long-horizon reasoning: 8-17 subgoals, 25-45 step episodes",
    "Partial observability: asymmetric info across heterogeneous agents",
    "Credit assignment: shaped reward over dependency-gated DAG",
    "Error recovery: invalid action returns error; agent must adapt",
]
limitations = [
    "Turn-based sequential - true parallel MARL needs action encoder",
    "PettingZoo AEC adapter is experimental, not fully tested",
    "Entity IDs fixed per task family (Factory V1) - not entity-disjoint",
    "Zero-shot measured on qwen2.5:3b only - model-dependent result",
    "No Figma, GitHub, or HR app simulators in this version",
]

add_rect(sl, 0.4, 1.85, 6.1, 4.8, fill=RGBColor(0xE6, 0xF6, 0xEF), line=0.8, line_color=GREEN_OK)
add_text(sl, "What this benchmark studies", 0.55, 1.9, 5.8, 0.42, size=13, bold=True, color=GREEN_OK)
for i, cap in enumerate(capabilities):
    add_text(sl, f"•  {cap}", 0.55, 2.4 + i*0.53, 5.8, 0.48, size=11, color=DARK_TEXT)

add_rect(sl, 6.8, 1.85, 6.1, 4.8, fill=RGBColor(0xFF, 0xF8, 0xE8), line=0.8, line_color=WARN)
add_text(sl, "Honest Limitations", 6.95, 1.9, 5.8, 0.42, size=13, bold=True, color=WARN)
for i, lim in enumerate(limitations):
    add_text(sl, f"•  {lim}", 6.95, 2.4 + i*0.72, 5.7, 0.62, size=11, color=DARK_TEXT)


# SLIDE 12 - RL Interface
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "RL & Research Interface",
            subtitle="Gymnasium-compatible · CTDE framing · evaluator state never exposed to policy")

add_rect(sl, 0.4, 1.85, 5.8, 2.9, fill=RGBColor(0xE6, 0xF0, 0xFF), line=0.8, line_color=ACCENT)
add_text(sl, "Observation (per-agent, role-filtered)", 0.55, 1.9, 5.5, 0.4,
         size=12, bold=True, color=ACCENT)
for i, line in enumerate([
    '"agent":    {employee_id, name, role, team_id}',
    '"inbox":    [email headers]  - body requires read_email',
    '"channels": [channel records]  - own channels only',
    '"calendar": [event records]  - own calendar only',
    '"sheets":   [{sheet_id, name, role}]',
    '"task":     {id, name, instruction}',
]):
    add_text(sl, line, 0.55, 2.38+i*0.38, 5.55, 0.35, size=9.5, color=DARK_TEXT)

add_rect(sl, 6.5, 1.85, 6.4, 1.35, fill=RGBColor(0xE6, 0xF6, 0xEF), line=0.8, line_color=GREEN_OK)
add_text(sl, "Action (structured semantic space)", 6.65, 1.9, 6.1, 0.4,
         size=12, bold=True, color=GREEN_OK)
add_text(sl,
    'Action(agent_id="pm_01", app="jira",\n'
    '       action_type="change_status",\n'
    '       parameters={"issue_id": "VEND-401", "status": "approved"})',
    6.6, 2.35, 6.15, 0.75, size=9.5, color=DARK_TEXT)

add_rect(sl, 6.5, 3.35, 6.4, 1.35, fill=RGBColor(0xFF, 0xF0, 0xD8), line=0.8, line_color=WARN)
add_text(sl, "Step return", 6.65, 3.4, 6.1, 0.4, size=12, bold=True, color=WARN)
add_text(sl,
    "obs, reward, terminated, truncated, info = env.step(action)\n\n"
    'info["eval"]  →  {progress, subgoals, reward_components}\n'
    "               # evaluator-only - NOT shown to agent",
    6.6, 3.8, 6.15, 0.75, size=9.5, color=DARK_TEXT)

add_rect(sl, 0.4, 4.95, 12.5, 1.45, fill=DARK_BG)
add_text(sl, "CTDE Framing", 0.6, 5.0, 3.0, 0.45, size=13, bold=True, color=ACCENT)
add_text(sl,
    'Training centralises info["eval"] subgoal labels for critic conditioning.\n'
    "Execution uses per-agent obs only - policies are fully decentralised.\n"
    "18 action types across 5 apps. Each agent sees only their permitted subset (14-18 actions).",
    0.6, 5.48, 12.2, 0.85, size=12, color=RGBColor(0xC0, 0xD8, 0xFF))


# SLIDE 13 - Research Questions
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 1: THE ENVIRONMENT")
slide_title(sl, "Research Questions Enabled",
            subtitle="The 80-point gap between zero-shot (0/30) and hint-guided (24/30) is the open problem")

rqs = [
    ("Can BC close 0% to 80%?",        "Behavioral Cloning on oracle trajectories",   "Trajectory export, deterministic seed"),
    ("CTDE vs IQL?",                    "MAPPO vs IQL (shared vs independent critic)", 'info["eval"]["subgoals"] labels'),
    ("Does action masking help?",       "PPO with/without permission mask",            "Static permission table per agent"),
    ("Distractor density effect?",      "RL across easy to adversarial presets",       "4 difficulty levels, calibrated counts"),
    ("Can policy generalise seeds?",    "Train 0-999k, test 2M-3M",                   "Disjoint seed-range splits"),
    ("Shaped vs sparse reward?",        "Ablate subgoal bonus",                        "Reward component breakdown"),
    ("Offline RL near-oracle?",         "IQL / CQL on oracle dataset",                "export_dataset()"),
]
col_x4 = [0.4, 4.5, 8.9]
col_w4 = [4.1, 4.4, 4.1]
add_rect(sl, 0.4, 1.85, 12.5, 0.45, fill=DARK_BG)
for hdr, cx, cw in zip(["Research Question", "Method", "Benchmark Feature"], col_x4, col_w4):
    add_text(sl, hdr, cx+0.05, 1.9, cw, 0.32, size=11, bold=True, color=WHITE)

for ri, row in enumerate(rqs):
    y = 2.3 + ri*0.52
    bg = CARD_BG if ri%2==0 else RGBColor(0xEF, 0xF3, 0xF7)
    add_rect(sl, 0.4, y, 12.5, 0.5, fill=bg, line=0.3, line_color=RGBColor(0xD0, 0xD8, 0xE0))
    for val, cx, cw in zip(row, col_x4, col_w4):
        add_text(sl, val, cx+0.05, y+0.1, cw, 0.32, size=10, color=DARK_TEXT)

add_rect(sl, 0.4, 6.6, 12.5, 0.65, fill=DARK_BG)
add_text(sl,
    "Next: Behavioral Cloning on hint-guided trajectories → PPO against shaped reward"
    " → entity-disjoint OOD evaluation via Factory V2",
    0.6, 6.7, 12.2, 0.45, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# SLIDE 14 - Section 2 transition
sl = prs.slides.add_slide(BLANK)
dark_slide(sl)
add_rect(sl, 0, 0, 13.33, 0.7, fill=ACCENT)
add_text(sl, "SECTION 2", 0.5, 0.12, 3.0, 0.45, size=14, bold=True, color=WHITE)

add_text(sl, "The Scenario Factory", 1.0, 2.3, 11.3, 1.2,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl,
    "A handcrafted task is a demo.\n"
    "A factory that mass-produces environments is a research instrument.",
    1.5, 3.7, 10.3, 0.9,
    size=18, color=RGBColor(0xB0, 0xD8, 0xFF), align=PP_ALIGN.CENTER, italic=True)

chips2 = [
    ("ScenarioFactory V1\nAll 6 tasks - IMPLEMENTED",          ACCENT2),
    ("Factory V2 - entity-level\nVendor onboarding - PROTOTYPE", WARN),
]
for i, (lbl, col) in enumerate(chips2):
    x = 2.0 + i*5.5
    add_rect(sl, x, 5.0, 4.5, 1.1, fill=col)
    add_text(sl, lbl, x+0.1, 5.08, 4.3, 0.95, size=13, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)


# SLIDE 15 - Why Factory Is Necessary
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 2: THE SCENARIO FACTORY")
slide_title(sl, "Why a Factory Is Necessary",
            subtitle="RL training requires seed-disjoint episodes - 6 handcrafted tasks are not enough")

needs = [
    ("Training data",   "Same 6 episodes every run",       "1,000+ seed-disjoint episodes"),
    ("Generalization",  "Impossible - train = test",       "Held-out seed spaces enforced"),
    ("Curriculum",      "Manual difficulty adjustment",    "easy to adversarial presets"),
    ("Reproducibility", "Hope nothing changed",            "Same seed, identical episode"),
    ("Scale-up",        "6 tasks maximum",                 "Thousands of validated instances"),
]
col_x5 = [0.5, 3.5, 8.2]
col_w5 = [3.0, 4.7, 4.7]
add_rect(sl, 0.5, 2.0, 12.3, 0.5, fill=DARK_BG)
for hdr, cx, cw in zip(["Need", "Without Factory", "With Factory"], col_x5, col_w5):
    add_text(sl, hdr, cx+0.1, 2.07, cw, 0.35, size=12, bold=True, color=WHITE)

for ri, row in enumerate(needs):
    y = 2.5 + ri*0.72
    bg = CARD_BG if ri%2==0 else RGBColor(0xEF, 0xF3, 0xF7)
    add_rect(sl, 0.5, y, 12.3, 0.68, fill=bg, line=0.4, line_color=RGBColor(0xD0, 0xD8, 0xE0))
    add_text(sl, row[0], 0.65, y+0.14, 2.8, 0.42, size=12, bold=True, color=DARK_TEXT)
    add_text(sl, row[1], 3.65, y+0.14, 4.4, 0.42, size=11, color=RED_NO)
    add_text(sl, row[2], 8.35, y+0.14, 4.4, 0.42, size=11, color=GREEN_OK, bold=True)


# SLIDE 16 - ScenarioFactory V1
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 2: THE SCENARIO FACTORY")
slide_title(sl, "ScenarioFactory V1 - Implemented",
            subtitle="All 6 tasks · 4 difficulty presets · seed-disjoint splits · validated manifests")

v1_caps = [
    ("build(task, seed, difficulty)",  "Live, validated env in one call"),
    ("4 difficulty presets",           "easy (2) / medium (6) / hard (15) / adversarial (30) distractors"),
    ("Realistic distractor pools",     "10 email subjects, 7 Jira tickets, 6 Slack messages"),
    ("Seed-disjoint splits",           "train: 0-999k  /  dev: 1M-2M  /  test: 2M-3M"),
    ("JSONL manifest export",          "scenario_id, seed, difficulty, apps, agents, subgoal_count"),
    ("Validation gates",               "Zero initial progress · acyclic DAG · dep IDs resolve"),
    ("6 task families",                "Full DAG, verifier, oracle baseline for every task"),
]
for i, (cap, detail) in enumerate(v1_caps):
    y = 1.9 + i*0.62
    add_rect(sl, 0.5, y, 12.3, 0.55,
             fill=CARD_BG if i%2==0 else RGBColor(0xEF, 0xF3, 0xF7),
             line=0.4, line_color=RGBColor(0xD0, 0xD8, 0xE0))
    add_rect(sl, 0.5, y, 0.35, 0.55, fill=ACCENT2)
    add_text(sl, "Yes", 0.5, y+0.1, 0.35, 0.35, size=9, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, cap, 0.95, y+0.1, 3.8, 0.38, size=11, bold=True, color=DARK_TEXT)
    add_text(sl, detail, 4.85, y+0.1, 7.9, 0.38, size=11, color=MID_GREY)

add_rect(sl, 0.5, 6.3, 12.3, 0.95, fill=DARK_BG)
add_text(sl,
    "factory = ScenarioFactory()\n"
    'env, obs, info = factory.build("vendor_onboarding", seed=42, difficulty="hard")\n'
    '# info["validator_status"] == "passed"',
    0.7, 6.35, 12.0, 0.82, size=10, color=RGBColor(0xC0, 0xE8, 0xFF))


# SLIDE 17 - Factory V2
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 2: THE SCENARIO FACTORY")
slide_title(sl, "Factory V2 - Entity-Level World Generation",
            subtitle="Vertical slice proof-of-concept on vendor_onboarding · same pipeline designed to extend")

add_rect(sl, 0.5, 1.85, 5.9, 4.5, fill=RGBColor(0xE6, 0xF0, 0xFF), line=0.8, line_color=ACCENT)
add_text(sl, "Seed 42", 0.65, 1.9, 5.5, 0.42, size=13, bold=True, color=ACCENT)
for i, line in enumerate([
    "Company: Smart Metrics",
    "Tickets: METR-401, METR-402, METR-403",
    "Unique employee names + emails",
    "sha256:841edac7c2667c37",
    "",
    "Oracle: PASS",
]):
    tc = GREEN_OK if "PASS" in line else DARK_TEXT
    add_text(sl, line, 0.65, 2.42+i*0.52, 5.5, 0.45, size=12, color=tc, bold=("PASS" in line))

add_rect(sl, 6.9, 1.85, 5.9, 4.5, fill=RGBColor(0xE8, 0xF6, 0xEE), line=0.8, line_color=ACCENT2)
add_text(sl, "Seed 43", 7.05, 1.9, 5.5, 0.42, size=13, bold=True, color=ACCENT2)
for i, line in enumerate([
    "Company: Stellar Services",
    "Tickets: SERV-401, SERV-402, SERV-403",
    "Different employee names + emails",
    "sha256:cbebc420d7920829",
    "",
    "Oracle: PASS",
]):
    tc = GREEN_OK if "PASS" in line else DARK_TEXT
    add_text(sl, line, 7.05, 2.42+i*0.52, 5.5, 0.45, size=12, color=tc, bold=("PASS" in line))

add_rect(sl, 0.5, 6.55, 12.3, 0.72, fill=DARK_BG)
add_text(sl,
    "Every generated world is oracle-verified at generation time. "
    "Same DAG structure, completely different operational context per seed.\n"
    "CompanySpec → WorldGenerator (isolated RNG) → WorldValidator → GeneratedVendorOnboardingTask",
    0.7, 6.6, 12.0, 0.6, size=11, color=RGBColor(0xB0, 0xD8, 0xFF))


# SLIDE 18 - Feature Status
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 2: THE SCENARIO FACTORY")
slide_title(sl, "Feature Status",
            subtitle="Delivered · Experimental · Planned - clearly separated")

cols_hdr  = ["IMPLEMENTED", "PROTOTYPE / EXPERIMENTAL", "ROADMAP"]
col_colors= [GREEN_OK, WARN, MID_GREY]
col_x6    = [0.4, 4.7, 9.0]
col_w6    = [4.3, 4.3, 4.2]

impl = [
    "6 tasks + DAG verifiers",
    "ScenarioFactory (all 6 tasks)",
    "Oracle + hint-guided + zero-shot baselines",
    "Deterministic DB-state verification",
    "Role-based permissions (RBAC)",
    "Shaped reward + anti-hacking logic",
    "Trajectory export + replay",
    "Streamlit dashboard",
    "Docker Compose deployment",
]
proto = [
    "Factory V2 world generation (1 task)",
    "PettingZoo AEC adapter",
    "SHA-256 world fingerprinting",
    "WorldValidator structural checks",
]
road = [
    "Factory V2 for all 6 tasks",
    "DAG synthesizer",
    "Org/identity generator",
    "Automated verifier generation",
    "Entity-disjoint OOD splits",
    "Large-scale curriculum calibration",
]

for ci, (hdr, color, cx, cw, items) in enumerate(
        zip(cols_hdr, col_colors, col_x6, col_w6, [impl, proto, road])):
    add_rect(sl, cx, 1.9, cw-0.1, 0.55, fill=color)
    add_text(sl, hdr, cx+0.1, 1.97, cw-0.2, 0.4,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, item in enumerate(items):
        add_rect(sl, cx, 2.5+ri*0.54, cw-0.1, 0.5,
                 fill=CARD_BG if ri%2==0 else RGBColor(0xEF, 0xF3, 0xF7),
                 line=0.3, line_color=RGBColor(0xD0, 0xD8, 0xE0))
        add_text(sl, item, cx+0.15, 2.57+ri*0.54, cw-0.3, 0.38, size=10, color=DARK_TEXT)

add_rect(sl, 0.4, 6.7, 12.5, 0.6, fill=DARK_BG)
add_text(sl,
    "118 tests pass (70 environment + 48 factory_v2)  ·  0 failed  ·  0 flaky",
    0.6, 6.78, 12.2, 0.4, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# SLIDE 19 - Difficulty & Reproducibility
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GREY)
section_bar(sl, "SECTION 2: THE SCENARIO FACTORY")
slide_title(sl, "Difficulty, Curriculum & Reproducibility",
            subtitle="4 distractor presets · seed-disjoint splits · same seed = identical episode always")

diff_rows = [
    ("easy",        "2",  "Minimal noise - near-direct path",         RGBColor(0x1A, 0xB0, 0x70)),
    ("medium",      "6",  "Focused retrieval required",               ACCENT),
    ("hard",        "15", "Heavy noise - same domain, wrong details", WARN),
    ("adversarial", "30", "Near-miss keywords designed to mislead",   RED_NO),
]
add_rect(sl, 0.5, 2.0, 9.0, 0.5, fill=DARK_BG)
for hdr, cx, cw in zip(["Preset", "Distractors", "Effect"], [0.6, 2.3, 4.0], [1.7, 1.7, 7.0]):
    add_text(sl, hdr, cx, 2.08, cw, 0.35, size=11, bold=True, color=WHITE)

for ri, (preset, n, effect, color) in enumerate(diff_rows):
    y = 2.5 + ri*0.65
    bg = CARD_BG if ri%2==0 else RGBColor(0xEF, 0xF3, 0xF7)
    add_rect(sl, 0.5, y, 9.0, 0.6, fill=bg, line=0.4, line_color=RGBColor(0xD0, 0xD8, 0xE0))
    add_rect(sl, 0.5, y, 0.3, 0.6, fill=color)
    add_text(sl, preset, 0.85, y+0.12, 1.4, 0.38, size=12, bold=True, color=color)
    add_text(sl, n,      2.4,  y+0.12, 1.5, 0.38, size=13, bold=True,
             color=DARK_TEXT, align=PP_ALIGN.CENTER)
    add_text(sl, effect, 4.1,  y+0.12, 5.3, 0.38, size=11, color=DARK_TEXT)

add_rect(sl, 9.7, 2.0, 3.3, 2.6, fill=DARK_BG)
add_text(sl, "Seed Splits", 9.85, 2.05, 3.0, 0.42, size=13, bold=True, color=ACCENT)
for i, (k, v) in enumerate([("train:", "0 to 999k"), ("dev:", "1M to 2M"), ("test:", "2M to 3M")]):
    add_text(sl, k, 9.85, 2.55+i*0.55, 1.2, 0.42, size=12, color=MID_GREY, bold=True)
    add_text(sl, v, 11.1, 2.55+i*0.55, 1.8, 0.42, size=12, color=WHITE)

add_rect(sl, 0.5, 5.25, 12.3, 1.05, fill=RGBColor(0xE8, 0xF4, 0xFF), line=0.8, line_color=ACCENT)
add_text(sl, "python scripts/generate_dataset.py \\", 0.7, 5.3, 12.0, 0.3, size=11, color=DARK_TEXT)
add_text(sl, "    --output generated_scenarios --train 1000 --dev 200 --test 500 --difficulty hard",
         0.7, 5.58, 12.0, 0.3, size=11, color=DARK_TEXT)
add_text(sl, "# Oracle = 100% at all difficulty levels  ·  Zero-Shot = 0% at all difficulty levels",
         0.7, 5.88, 12.0, 0.3, size=10, color=MID_GREY, italic=True)

add_rect(sl, 0.5, 6.5, 12.3, 0.72, fill=DARK_BG)
add_text(sl,
    "Disjoint seed ranges prevent episode reuse across train/dev/test. "
    "Same seed → identical episode, always - fully reproducible.",
    0.7, 6.6, 12.0, 0.5, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# SLIDE 20 - Closing
sl = prs.slides.add_slide(BLANK)
dark_slide(sl)
add_rect(sl, 0, 0, 13.33, 0.55, fill=ACCENT)
add_text(sl, "SECTION 2: THE SCENARIO FACTORY  ·  CLOSING", 0.3, 0.08, 10, 0.4,
         size=12, bold=True, color=WHITE)

add_text(sl, "The Open Research Question", 0.5, 0.65, 12.3, 0.65,
         size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, 1.5, 1.5, 2.8, 1.0, fill=MID_GREY)
add_text(sl, "Zero-Shot LLM", 1.5, 1.52, 2.8, 0.35, size=11, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "0 / 30", 1.5, 1.85, 2.8, 0.55, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "← 80-point gap →", 4.4, 1.8, 4.5, 0.55, size=22, bold=True,
         color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(sl, 9.0, 1.5, 2.8, 1.0, fill=ACCENT)
add_text(sl, "Hint-Guided LLM", 9.0, 1.52, 2.8, 0.35, size=11, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "24 / 30", 9.0, 1.85, 2.8, 0.55, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl,
    "Can autonomous multi-agent systems close the gap between\n"
    "zero-shot performance (0/30) and procedurally-guided performance (24/30)\n"
    "in realistic enterprise coordination workflows?",
    1.0, 2.7, 11.3, 1.1, size=16, color=RGBColor(0xB0, 0xD8, 0xFF),
    align=PP_ALIGN.CENTER, italic=True)

deliv = [
    "6 long-horizon multi-agent tasks",
    "5 apps · 5 agents · role-based permissions",
    "Dependency-gated subgoal DAGs",
    "Deterministic DB-state verification",
    "Oracle 30/30 · Hint-guided 24/30 · Zero-Shot 0/30",
    "ScenarioFactory: seeds x difficulty x splits",
    "Factory V2: entity-level world generation (prototype)",
    "118 automated tests · 0 flaky",
    "Docker Compose reproducible deployment",
]
for i, item in enumerate(deliv):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = 0.5 + col*6.5
    y = 4.0 + row*0.47
    add_text(sl, f"  {item}", x, y, 6.2, 0.42, size=11, color=WHITE)

add_rect(sl, 0.5, 6.5, 12.3, 0.75, fill=RGBColor(0x1A, 0x3A, 0x5C))
add_text(sl,
    "Next: Behavioral Cloning → PPO against shaped reward → entity-disjoint OOD eval via Factory V2",
    0.7, 6.62, 12.0, 0.5, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# Save
import os
out_path = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..", "deliverables", "enterprise_marl_benchmark.pptx"))
prs.save(out_path)
print(f"Saved: {out_path}")
